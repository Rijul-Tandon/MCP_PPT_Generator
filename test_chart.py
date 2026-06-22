from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
slide_layout = prs.slide_layouts[8] # Title and Chart layout usually
slide = prs.slides.add_slide(slide_layout)

chart_data = CategoryChartData()
chart_data.categories = ['A', 'B', 'C']
chart_data.add_series('Series 1', (1, 2, 3))

for shape in slide.placeholders:
    if shape.is_placeholder:
        print("Found placeholder:", shape.placeholder_format.type)
        if shape.placeholder_format.type == 14: # CHART
            chart = shape.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
            print("Inserted chart!")
            break

prs.save("test_chart.pptx")
