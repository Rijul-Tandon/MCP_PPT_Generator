from pptx import Presentation

def replace_text_test():
    out_path = 'output/clone_test.pptx'
    prs = Presentation(out_path)
    
    # We mapped [1, 5, 5, 10] (1-based), so slide index 1 (0-based) is the 5th slide of original deck
    # That slide is "Leveraging the referral network, we deep-dive into insights across geography, patient flow, and influence dynamics"
    slide = prs.slides[1]
    
    replacements = {
        "Leveraging the referral network": "Leveraging the amazing LLM network"
    }
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        for old_t, new_t in replacements.items():
            if old_t in text:
                print(f"Found '{old_t}' in shape '{shape.name}'")
                # To preserve styling, we replace text in the paragraphs/runs
                # A simple replace: just put all text in the first run and clear the rest
                new_full_text = text.replace(old_t, new_t)
                p = shape.text_frame.paragraphs[0]
                r = p.runs[0] if p.runs else p.add_run()
                r.text = new_full_text
                # Clear other paragraphs/runs to avoid duplication
                # (This is a simplified test, a robust version would be more careful)
                for i in range(1, len(shape.text_frame.paragraphs)):
                    p = shape.text_frame.paragraphs[i]
                    p.clear() # clear runs
    
    prs.save('output/clone_test_replaced.pptx')
    print("Replaced text successfully")

if __name__ == '__main__':
    replace_text_test()
