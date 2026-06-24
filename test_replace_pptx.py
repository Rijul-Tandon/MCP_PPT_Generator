from pptx import Presentation
import os
import shutil

def test_replace():
    src = os.path.abspath('context/Referral_Analysis/reference_decks/2026-05-14 NT1 Referral Network Analysis - Insights & Recommendations v1.pptx')
    dst = os.path.abspath('output/test_replace.pptx')
    shutil.copy(src, dst)
    
    prs = Presentation(dst)
    slide = prs.slides[0]
    
    mapping = {
        "TAK-861 HCP Referral Analysis Insights & Recommendations": "LLM Generated Agenda Slide",
        "Focus of the conversation today": "Agenda Overview"
    }
    
    replaced = 0
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for old_txt, new_txt in mapping.items():
            if old_txt in shape.text:
                # Safe replace algorithm
                for p in shape.text_frame.paragraphs:
                    if old_txt in p.text:
                        if not p.runs: continue
                        # Assign replaced text to first run
                        new_p_text = p.text.replace(old_txt, new_txt)
                        p.runs[0].text = new_p_text
                        # Clear other runs
                        for r in p.runs[1:]:
                            r.text = ""
                        replaced += 1
                        print(f"Replaced '{old_txt}' -> '{new_txt}'")
                        
    prs.save(dst)
    print("Done, replaced", replaced)

if __name__ == "__main__":
    test_replace()
