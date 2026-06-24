from pptx import Presentation
import os

def test_extract():
    path = os.path.abspath('context/referral_analysis/reference_decks/2026-05-14 NT1 Referral Network Analysis - Insights & Recommendations v1.pptx')
    prs = Presentation(path)
    
    slide = prs.slides[0]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                texts.append(text)
                
    print("Extracted Texts from Slide 0:")
    for t in texts:
        print(f" - {repr(t)}")

if __name__ == "__main__":
    test_extract()
