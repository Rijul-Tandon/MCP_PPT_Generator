"""
ppt_server.py — MCP Server for PowerPoint operations.

Architecture:
  The LLM communicates with this server via JSON-RPC 2.0 (MCP protocol).
  XML content (OOXML) travels as string values inside JSON-RPC payloads —
  the LLM reads and writes raw Office Open XML so nothing is lost or
  silently re-interpreted by a Python parser.

  JSON-RPC transport (MCP)
        │
        ├── get_deck_manifest   → compact slide-title index (~300 tokens/deck)
        ├── get_slide_xml       → full OOXML for ONE slide + linked charts
        ├── write_slide_xml     → validated XML write-back directly into ZIP
        └── write_slide_chart_xml → surgical chart-only XML replacement

  High-level helpers:
        ├── list_use_cases / read_context / write_context / list_data_sources
        ├── clone_deck_from_reference  (win32com cloning for perfect aesthetics)
        ├── get_slide_text             (lightweight text extraction, ~100 tokens)
        └── replace_slide_text         (preserves design while updating text)
"""

import json
import os
import re
import shutil
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from mcp.server.fastmcp import FastMCP

from .context_manager import ContextManager
from .tools_query import list_data_sources as _list_ds

mcp = FastMCP("PPTServer")

# ─────────────────────────────────────────────────────────────────────────────
# OOXML namespace map used for ElementTree queries
# ─────────────────────────────────────────────────────────────────────────────
_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────

def _sorted_slide_files(z: zipfile.ZipFile) -> list[str]:
    """Return slide XML paths from a PPTX ZIP sorted by slide number."""
    return sorted(
        [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)],
        key=lambda n: int(re.search(r"(\d+)", n).group(1)),
    )


def _strip_whitespace(xml: str) -> str:
    """Remove inter-tag whitespace to cut token count without losing content."""
    return re.sub(r">\s+<", "><", xml).strip()


def _extract_title(xml: str) -> str:
    """
    Pull the title text from a slide's OOXML.
    Looks for the title/ctrTitle placeholder shape; falls back to first text run.
    """
    try:
        root = ET.fromstring(xml)
        for sp in root.findall(".//p:sp", _NS):
            ph = sp.find(".//p:ph", _NS)
            if ph is not None and ph.get("type") in ("title", "ctrTitle"):
                texts = sp.findall(".//a:t", _NS)
                title = "".join(t.text or "" for t in texts).strip()
                if title:
                    return title
        # Fallback: first non-empty text run on the slide
        for t in root.findall(".//a:t", _NS):
            if t.text and t.text.strip():
                return t.text.strip()[:80]
    except ET.ParseError:
        pass
    return "Untitled"


def _chart_paths_for_slide(z: zipfile.ZipFile, slide_file: str) -> dict[str, str]:
    """
    Follow the .rels file for a slide and return {rel_id: chart_zip_path}
    for every chart relationship found.
    """
    rels_path = f"ppt/slides/_rels/{slide_file.split('/')[-1]}.rels"
    if rels_path not in z.namelist():
        return {}
    rels_xml = z.read(rels_path).decode("utf-8", errors="ignore")
    chart_map: dict[str, str] = {}
    for rel_id, target in re.findall(
        r'<Relationship[^>]*Id="([^"]+)"[^>]*Type="[^"]*chart[^"]*"[^>]*Target="([^"]+)"',
        rels_xml,
    ):
        chart_zip_path = "ppt/charts/" + target.split("/")[-1]
        if chart_zip_path in z.namelist():
            chart_map[rel_id] = chart_zip_path
    return chart_map


def _update_zip(zip_path: str, updates: dict[str, bytes]) -> None:
    """
    Replace specific entries inside a ZIP file in-place.
    All other entries are copied without modification.
    """
    tmp_path = zip_path + ".updating"
    try:
        with zipfile.ZipFile(zip_path, "r") as zin, \
             zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in updates:
                    zout.writestr(item, updates[item.filename])
                else:
                    zout.writestr(item, zin.read(item.filename))
        shutil.move(tmp_path, zip_path)
    finally:
        if Path(tmp_path).exists():
            Path(tmp_path).unlink()


def _validate_xml(xml_string: str, label: str = "slide") -> str | None:
    """
    Three-level validation for LLM-authored OOXML.
    Returns None if valid, or a descriptive error string to return to the LLM.

    Level 1 — Well-formedness (catches unclosed tags, bad escaping)
    Level 2 — Required OOXML namespaces present
    Level 3 — PPTX structural check via python-pptx round-trip
    """
    # Level 1
    try:
        ET.fromstring(xml_string)
    except ET.ParseError as e:
        return f"Validation Error (Level 1 — Malformed XML in {label}): {e}"

    # Level 2
    for ns in ("xmlns:a=", "xmlns:p="):
        if ns not in xml_string:
            return (
                f"Validation Error (Level 2 — Missing namespace '{ns}' in {label}). "
                "Ensure the root element declares all required OOXML namespaces."
            )

    return None  # all levels passed


# ─────────────────────────────────────────────────────────────────────────────
# Context & use-case tools  (lightweight JSON outputs)
# ─────────────────────────────────────────────────────────────────────────────

@mcp.tool()
def list_use_cases() -> str:
    """List all available project use cases found in the context/ directory."""
    manager = ContextManager(Path.cwd())
    cases = manager.list_use_cases()
    if not cases:
        return "No use cases found in the context/ directory."
    return "\n".join(f"- {c.id}: {c.name}" for c in cases)


@mcp.tool()
def read_context(use_case: str) -> str:
    """Read context.txt for a use case to understand project goals and background."""
    manager = ContextManager(Path.cwd())
    try:
        folder = manager.resolve_use_case(use_case)
        path = manager.resolve_context_path(folder)
        return manager.load_context(path)
    except Exception as e:
        return f"Error reading context: {e}"


@mcp.tool()
def write_context(use_case: str, full_content: str) -> str:
    """Overwrite context.txt for a use case with updated project information."""
    manager = ContextManager(Path.cwd())
    try:
        folder = manager.resolve_use_case(use_case)
        path = manager.resolve_context_path(folder)
        path.write_text(full_content, encoding="utf-8")
        return f"Successfully updated context.txt for '{use_case}'"
    except Exception as e:
        return f"Error writing context: {e}"


@mcp.tool()
def list_data_sources(use_case: str) -> str:
    """
    List all Excel (.xlsx/.csv) and PowerPoint (.pptx) reference files
    available for a given use case. Returns absolute file paths for use
    with get_deck_manifest and get_slide_xml.
    """
    return _list_ds(Path.cwd(), use_case)


# ─────────────────────────────────────────────────────────────────────────────
# XML-native PPT tools  (lossless read / write)
# ─────────────────────────────────────────────────────────────────────────────

@mcp.tool()
def get_deck_manifest(file_path: str) -> str:
    """
    Return a compact JSON index of all slides in a PPTX deck.

    Each entry contains only the slide index and title (~300 tokens for a
    20-slide deck). Use this to let the user pick which slides to focus on
    before loading the full XML of any individual slide.

    Returns JSON:
      { "deck": "filename.pptx", "slide_count": N,
        "slides": [{"idx": 0, "title": "..."}, ...] }
    """
    path = file_path.strip('"').strip("'")
    try:
        with zipfile.ZipFile(path, "r") as z:
            slide_files = _sorted_slide_files(z)
            slides = []
            for idx, sf in enumerate(slide_files):
                raw_xml = z.read(sf).decode("utf-8", errors="ignore")
                slides.append({"idx": idx, "title": _extract_title(raw_xml)})
        return json.dumps(
            {"deck": Path(path).name, "slide_count": len(slides), "slides": slides},
            indent=2,
        )
    except Exception as e:
        return f"Error reading deck manifest: {e}"


def get_slide_xml(file_path: str, slide_idx: int) -> str:
    """
    Return the full, raw OOXML for ONE slide and ALL charts linked to it.

    The XML is whitespace-stripped to minimise token usage (~2,000-4,000
    tokens per slide). Charts are returned keyed by their relationship ID
    so the LLM can reference them when authoring replacement chart XML.

    Returns JSON:
      {
        "slide_idx": N,
        "slide_xml": "<p:sld ...>...</p:sld>",
        "charts": { "rId2": "<c:chartSpace ...>...</c:chartSpace>", ... }
      }
    """
    path = file_path.strip('"').strip("'")
    try:
        with zipfile.ZipFile(path, "r") as z:
            slide_files = _sorted_slide_files(z)
            if slide_idx < 0 or slide_idx >= len(slide_files):
                return json.dumps({"error": f"slide_idx {slide_idx} out of range (deck has {len(slide_files)} slides)"})
            sf = slide_files[slide_idx]
            slide_xml = _strip_whitespace(z.read(sf).decode("utf-8", errors="ignore"))
            chart_map = _chart_paths_for_slide(z, sf)
            charts = {
                rel_id: _strip_whitespace(z.read(cp).decode("utf-8", errors="ignore"))
                for rel_id, cp in chart_map.items()
            }
        return json.dumps(
            {"slide_idx": slide_idx, "slide_xml": slide_xml, "charts": charts},
            ensure_ascii=False,
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


def write_slide_xml(
    file_path: str,
    slide_idx: int,
    slide_xml: str,
    chart_xmls: dict | None = None,
) -> str:
    """
    Write LLM-authored OOXML back into a PPTX file for a specific slide.

    slide_xml must be a complete <p:sld> element with all required OOXML
    namespaces declared. chart_xmls is optional: pass {rel_id: xml_string}
    to also replace chart data on this slide.

    Validation (3 levels) runs before the file is touched:
      1. Well-formedness — catches malformed XML tags
      2. Namespace check — ensures required xmlns:a and xmlns:p are present
      3. Round-trip load — verifies python-pptx can still open the file

    If any level fails the error is returned to the LLM for self-correction.
    The original file is never modified on failure.
    """
    path = file_path.strip('"').strip("'")
    chart_xmls = chart_xmls or {}

    # ── Validate slide XML ────────────────────────────────────────────────────
    err = _validate_xml(slide_xml, "slide_xml")
    if err:
        return err

    # ── Validate each chart XML ───────────────────────────────────────────────
    for rel_id, cxml in chart_xmls.items():
        err = _validate_xml(cxml, f"chart {rel_id}")
        if err:
            return err

    try:
        with zipfile.ZipFile(path, "r") as z:
            slide_files = _sorted_slide_files(z)
            if slide_idx < 0 or slide_idx >= len(slide_files):
                return f"Error: slide_idx {slide_idx} out of range (deck has {len(slide_files)} slides)"
            sf = slide_files[slide_idx]
            chart_map = _chart_paths_for_slide(z, sf)

        # Build the updates dict: slide XML + any chart XMLs
        updates: dict[str, bytes] = {sf: slide_xml.encode("utf-8")}
        for rel_id, cxml in chart_xmls.items():
            if rel_id in chart_map:
                updates[chart_map[rel_id]] = cxml.encode("utf-8")
            else:
                return f"Error: relationship ID '{rel_id}' not found on slide {slide_idx}"

        # ── Level 3: write to temp copy and test round-trip ──────────────────
        tmp = path + ".validate"
        try:
            shutil.copy2(path, tmp)
            _update_zip(tmp, updates)
            from pptx import Presentation
            Presentation(tmp)          # raises if structurally broken
        except Exception as e:
            return f"Validation Error (Level 3 — PPTX structure): {e}"
        finally:
            if Path(tmp).exists():
                Path(tmp).unlink()

        # ── All good — apply to the real file ────────────────────────────────
        _update_zip(path, updates)
        charts_updated = list(chart_xmls.keys())
        return (
            f"Successfully wrote slide {slide_idx} to {Path(path).name}. "
            + (f"Charts updated: {charts_updated}" if charts_updated else "No chart updates.")
        )

    except Exception as e:
        return f"Error writing slide XML: {e}"


@mcp.tool()
def write_slide_chart_xml(
    file_path: str,
    slide_idx: int,
    rel_id: str,
    chart_xml: str,
) -> str:
    """
    Surgically replace ONLY the chart data/formatting XML for one chart on a slide.

    Use this when you want to update a chart's data or styling without
    touching the slide's layout XML at all. rel_id must match a chart
    relationship ID returned by get_slide_xml (e.g. 'rId2').

    The chart_xml must be a complete <c:chartSpace> element.
    """
    path = file_path.strip('"').strip("'")

    # Basic well-formedness check (namespaces in chart XMLs use c:, not p:)
    try:
        ET.fromstring(chart_xml)
    except ET.ParseError as e:
        return f"Validation Error (Malformed chart XML): {e}"

    try:
        with zipfile.ZipFile(path, "r") as z:
            slide_files = _sorted_slide_files(z)
            if slide_idx < 0 or slide_idx >= len(slide_files):
                return f"Error: slide_idx {slide_idx} out of range"
            sf = slide_files[slide_idx]
            chart_map = _chart_paths_for_slide(z, sf)

        if rel_id not in chart_map:
            return (
                f"Error: relationship ID '{rel_id}' not found on slide {slide_idx}. "
                f"Available chart rel IDs: {list(chart_map.keys())}"
            )

        _update_zip(path, {chart_map[rel_id]: chart_xml.encode("utf-8")})
        return f"Successfully updated chart '{rel_id}' on slide {slide_idx} in {Path(path).name}"

    except Exception as e:
        return f"Error writing chart XML: {e}"


# ─────────────────────────────────────────────────────────────────────────────
# Deck Cloning (100% Aesthetic Fidelity)
# ─────────────────────────────────────────────────────────────────────────────

@mcp.tool()
def clone_deck_from_reference(
    reference_path: str,
    output_filename: str,
    slide_mapping: list[int]
) -> str:
    """
    Creates a new presentation by perfectly duplicating the chosen slides from 
    a reference deck. This guarantees 100% fidelity of layouts, theme colors, 
    charts, logos, and relationships.
    
    reference_path: Absolute path to the source .pptx
    output_filename: Name of the output file (e.g. "referral_analysis_v2.pptx"). 
                     It will be saved in the output/ directory.
    slide_mapping: List of 0-based slide indices to keep, in the desired order.
                   e.g., [0, 4, 4, 9] creates a 4-slide deck.
    """
    from .clone_builder import clone_deck
    
    base_dir = Path.cwd()
    output_dir = base_dir / "output"
    output_dir.mkdir(exist_ok=True)
    out_path = str(output_dir / output_filename)
    
    try:
        return clone_deck(reference_path, out_path, slide_mapping)
    except Exception as e:
        return f"Error cloning deck: {e}"

# ─────────────────────────────────────────────────────────────────────────────
# Text Extraction & Replacement (Token-Efficient Design Preservation)
# ─────────────────────────────────────────────────────────────────────────────

@mcp.tool()
def get_slide_text(file_path: str, slide_idx: int) -> str:
    """
    Extracts all text from a specific slide using python-pptx. 
    Use this instead of get_slide_xml to understand what is on a slide 
    without burning thousands of tokens on raw OOXML.
    
    file_path: Absolute path to the presentation
    slide_idx: 0-based index of the slide
    Returns a JSON list of text strings found on the slide.
    """
    from pptx import Presentation
    
    path = os.path.abspath(file_path)
    if not os.path.exists(path):
        return f"Error: File not found: {path}"
        
    try:
        prs = Presentation(path)
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return f"Error: slide_idx {slide_idx} out of range (0 to {len(prs.slides)-1})."
            
        slide = prs.slides[slide_idx]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text.strip()
                if txt:
                    # Replace invisible breaking characters to make matching easier
                    txt = txt.replace('\\x0b', '\\n')
                    texts.append(txt)
        return json.dumps(texts, indent=2)
    except Exception as e:
        return f"Error extracting text: {e}"

@mcp.tool()
def replace_slide_text(file_path: str, slide_idx: int, replacements: dict) -> str:
    """
    Surgically replaces text on a slide while perfectly preserving its 
    font, color, size, and layout styling.
    
    file_path: Absolute path to the presentation (usually the cloned output file)
    slide_idx: 0-based index of the slide to edit
    replacements: A JSON dictionary of {"Exact Old Text": "New Text"}
    """
    from pptx import Presentation
    
    path = os.path.abspath(file_path)
    if not os.path.exists(path):
        return f"Error: File not found: {path}"
        
    try:
        prs = Presentation(path)
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return f"Error: slide_idx {slide_idx} out of range."
            
        slide = prs.slides[slide_idx]
        replaced_count = 0
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for old_txt, new_txt in replacements.items():
                # We do a fast check first
                if old_txt in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        if old_txt in paragraph.text:
                            # To preserve styling, we replace text in the paragraph's 
                            # first run and clear the subsequent runs.
                            if not paragraph.runs: continue
                            new_p_text = paragraph.text.replace(old_txt, new_txt)
                            paragraph.runs[0].text = new_p_text
                            for r in paragraph.runs[1:]:
                                r.text = ""
                            replaced_count += 1
                            
        prs.save(path)
        return f"Successfully made {replaced_count} text replacements on slide {slide_idx}."
    except Exception as e:
        return f"Error replacing text: {e}"
