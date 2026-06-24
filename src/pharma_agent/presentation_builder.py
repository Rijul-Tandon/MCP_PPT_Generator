from __future__ import annotations

import json # Used to write output trace logs in JSON format
from pathlib import Path # Used to resolve paths for templates and output locations


class PresentationBuilder:
    """Render slide draft JSON into a PowerPoint file.

    There are two write modes:
    1. Build mode: create a new deck from a template shell.
    2. Refine mode: keep the existing deck's slides and update them in place.

    Refine mode is intentionally conservative. It updates title and text-bearing
    regions first, preserves the original theme/layout/visual objects as much as
    possible, and only creates new slides when the new outline is longer than the
    existing deck.
    """

    def __init__(self) -> None:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.dml.color import RGBColor  # type: ignore
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER  # type: ignore
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for the build command. Install it with `pip install -r requirements.txt`."
            ) from exc

        self.Presentation = Presentation
        self.RGBColor = RGBColor
        self.MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE
        self.PP_PLACEHOLDER = PP_PLACEHOLDER
        self.MSO_ANCHOR = MSO_ANCHOR
        self.PP_ALIGN = PP_ALIGN
        self.Inches = Inches
        self.Pt = Pt

    def build_from_content(
        self,
        content: dict,
        pptx_path: Path,
        trace_path: Path | None = None,
        existing_presentation_path: Path | None = None,
    ) -> list[str]:
        """Execute the PPTX generation by updating an existing deck or creating a new one."""
        warnings: list[str] = []
        presentation = self._load_presentation(content, existing_presentation_path, warnings)
        trace = self._create_trace_shell(content)
        self._apply_metadata(presentation, content)

        if existing_presentation_path and existing_presentation_path.exists():
            self._refine_existing_presentation(presentation, content, trace, warnings)
        else:
            # Build mode clears the template shell because we want only the new story.
            self._reset_to_template_shell(presentation)
            for slide_data in content.get("slides", []):
                slide = self._create_slide(presentation, slide_data)
                self._record_slide_trace(trace, slide_data)
                self._warn_if_claims_missing(slide_data, warnings)
                self._attach_notes(slide, slide_data)

        pptx_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(pptx_path))
        if trace_path is not None:
            trace_path.write_text(json.dumps(trace, indent=2), encoding="utf-8")
        return warnings

    # ── Methods removed in XML-native refactor ─────────────────────────────
    # read_structure, update_slide, add_chart, add_table, add_floating_textbox
    # were replaced by ppt_server.get_slide_xml / write_slide_xml /
    # write_slide_chart_xml which work directly on OOXML inside the ZIP.
    # ────────────────────────────────────────────────────────────────────────


    def _load_presentation(self, content: dict, existing_presentation_path: Path | None, warnings: list[str]):
        """Load the base presentation from an existing file, a template, or start blank."""
        # Refinement mode starts from the existing deck so we can preserve its theme,
        # masters, and layout family.
        if existing_presentation_path and existing_presentation_path.exists():
            return self.Presentation(str(existing_presentation_path))

        template_path = Path(content.get("templateDeckPath", "")) if content.get("templateDeckPath") else None
        if template_path and template_path.exists():
            return self.Presentation(str(template_path))

        warnings.append("Template deck could not be resolved. Falling back to a blank presentation shell.")
        presentation = self.Presentation()
        presentation.slide_width = self.Inches(13.333)
        presentation.slide_height = self.Inches(7.5)
        return presentation

    def _refine_existing_presentation(self, presentation, content: dict, trace: dict, warnings: list[str]) -> None:
        """Update existing slides with newly planned content in-place."""
        existing_slides = list(presentation.slides)
        slide_payloads = content.get("slides", [])

        # Update matching slides in place so manual layouts and visuals survive.
        for index, slide_data in enumerate(slide_payloads):
            if index < len(existing_slides):
                slide = existing_slides[index]
                self._refine_slide(slide, slide_data)
            else:
                slide = self._create_slide(presentation, slide_data)
                warnings.append(f"Created new slide {slide_data.get('slideNumber')} because the existing deck was shorter than the new outline.")
            self._record_slide_trace(trace, slide_data)
            self._warn_if_claims_missing(slide_data, warnings)
            self._attach_notes(slide, slide_data)

        if len(existing_slides) > len(slide_payloads):
            warnings.append(
                f"Existing deck contains {len(existing_slides) - len(slide_payloads)} trailing slide(s) beyond the refined outline. They were left unchanged."
            )

    def _refine_slide(self, slide, slide_data: dict) -> None:
        """Apply outline content to a specific slide, preserving format where possible."""
        self._set_title(slide, slide_data.get("title", "Untitled"), top=0.65, size=24)

        blocks = self._build_refinement_blocks(slide_data)
        editable_shapes = self._find_editable_text_shapes(slide)

        # Reuse existing text regions first; only add textboxes if the old slide does
        # not have enough editable areas for the new content.
        for shape, block in zip(editable_shapes, blocks):
            self._set_shape_text(shape, block, 15)

        if len(blocks) > len(editable_shapes):
            self._add_missing_text_blocks(slide, blocks[len(editable_shapes):])

        self._upsert_footer(slide, slide_data)

    def _build_refinement_blocks(self, slide_data: dict) -> list[str]:
        blocks: list[str] = []
        objective = slide_data.get("objective", "").strip()
        bullets = slide_data.get("bullets", [])
        visual = slide_data.get("visualSuggestion", "").strip()
        if objective:
            blocks.append(objective)
        if bullets:
            blocks.append("\n".join(f"- {bullet}" for bullet in bullets))
        if visual:
            blocks.append(f"Recommended visual direction: {visual}")
        return blocks

    def _find_editable_text_shapes(self, slide) -> list:
        editable = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            if self._is_title_shape(shape):
                continue
            text = self._shape_text(shape)
            if text.strip().startswith("Claim sources:"):
                continue
            editable.append(shape)
        editable.sort(key=lambda shape: (getattr(shape, "top", 0), -(getattr(shape, "width", 0) * getattr(shape, "height", 0))))
        return editable

    def _add_missing_text_blocks(self, slide, blocks: list[str]) -> None:
        top = 1.7
        for block in blocks:
            self._add_textbox(slide, 0.95, top, 10.8, 0.9, block, 14)
            top += 1.1

    def _is_title_shape(self, shape) -> bool:
        try:
            return shape.placeholder_format.type == self.PP_PLACEHOLDER.TITLE
        except Exception:
            return False

    def _shape_text(self, shape) -> str:
        try:
            return shape.text_frame.text
        except Exception:
            return ""

    def _create_trace_shell(self, content: dict) -> dict:
        return {
            "title": content.get("title", "Presentation"),
            "templateDeckPath": str(content.get("templateDeckPath", "")),
            "slides": [],
        }

    def _record_slide_trace(self, trace: dict, slide_data: dict) -> None:
        trace["slides"].append(
            {
                "slideNumber": slide_data.get("slideNumber"),
                "title": slide_data.get("title"),
                "archetype": slide_data.get("archetype"),
                "layoutHint": slide_data.get("layoutHint"),
                "claimSources": slide_data.get("claimSources", []),
                "styleSources": slide_data.get("styleSources", []),
                "notes": slide_data.get("notes", []),
            }
        )

    def _warn_if_claims_missing(self, slide_data: dict, warnings: list[str]) -> None:
        if not slide_data.get("claimSources") and slide_data.get("archetype") not in {"agenda", "section_divider"}:
            warnings.append(f"Slide {slide_data.get('slideNumber')} has no claim-bearing sources.")

    def _reset_to_template_shell(self, presentation) -> None:
        for slide_id in list(presentation.slides._sldIdLst):
            presentation.part.drop_rel(slide_id.rId)
            presentation.slides._sldIdLst.remove(slide_id)

    def _apply_metadata(self, presentation, content: dict) -> None:
        presentation.core_properties.author = "Pharma Presentation Agent"
        presentation.core_properties.title = content.get("title", "Pharma Presentation Agent Output")
        presentation.core_properties.subject = "Business-development presentation draft"

    def _create_slide(self, presentation, slide_data: dict):
        layout = self._choose_layout(presentation, slide_data)
        slide = presentation.slides.add_slide(layout)
        self._clear_text(slide)
        self._render_slide(slide, slide_data)
        self._add_footer(slide, slide_data)
        return slide

    def _render_slide(self, slide, slide_data: dict) -> None:
        title_text = slide_data.get("title", "Untitled")
        bullets = slide_data.get("bullets", [])
        visual_hint = slide_data.get("visualSuggestion", "")

        title_set = False
        body_set = False
        visual_set = False

        # STRICTLY use master placeholders to preserve template design
        for shape in slide.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                
                if ph_type == self.PP_PLACEHOLDER.TITLE and not title_set:
                    shape.text = title_text
                    title_set = True
                    
                elif ph_type in (self.PP_PLACEHOLDER.BODY, self.PP_PLACEHOLDER.OBJECT) and not body_set and bullets:
                    tf = shape.text_frame
                    tf.clear()
                    for i, bullet in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                    body_set = True
                    
                elif ph_type in (self.PP_PLACEHOLDER.PICTURE, self.PP_PLACEHOLDER.CHART, self.PP_PLACEHOLDER.TABLE, self.PP_PLACEHOLDER.OBJECT) and not visual_set and visual_hint:
                    # If information/chart is missing, inject guided text into the visual placeholder
                    shape.text = f"[GUIDED VISUAL REQUIRED]\n\n{visual_hint}\n\n(Please replace this placeholder with the appropriate chart/visual from the template.)"
                    visual_set = True
            except Exception:
                continue

        # Fallbacks: If the layout didn't have enough placeholders, we gracefully inject
        # explicit fallback shapes with guided text so no content is lost.
        if not title_set:
            self._set_title(slide, title_text, top=0.5, size=24)
            
        if not body_set and bullets:
            self._add_text_block(slide, 1.0, 1.5, 5.0, 4.0, bullets, 16)
            
        if not visual_set and visual_hint:
            self._add_visual_placeholder(slide, 6.5, 1.5, 5.0, 4.0, visual_hint)

    def _choose_layout(self, presentation, slide_data: dict):
        preferred_names = {
            "title": ["Title Page 1", "Title Slide", "Title Page 2"],
            "agenda": ["Section Title Page", "2_Standard 1-Column Text", "Standard 1-Column Text"],
            "executive_summary": ["1_Exec Sum", "Exec Sum", "Exec Sum2", "2_Exec Sum"],
            "framework": ["Standard 2-Column Text", "2_Standard 1-Column Text", "Standard"],
            "patient_funnel": ["Patient_Flow", "Advanced Chart Full Width", "Standard 2-Column Text"],
            "insight": ["Advanced Chart Full Width", "Advanced Chart 2/3", "Standard 2-Column Text"],
            "recommendation": ["Standard 2-Column Text", "2_Standard 1-Column Text", "Exec Sum"],
            "appendix": ["2_Header Only", "Standard 1-Column Text", "Blank slide"],
        }
        for name in preferred_names.get(slide_data.get("archetype", "insight"), ["Standard 2-Column Text", "Blank slide"]):
            for layout in presentation.slide_layouts:
                if layout.name.strip().lower() == name.strip().lower():
                    return layout
        return presentation.slide_layouts[0 if len(presentation.slide_layouts) else 0]

    def _clear_text(self, slide) -> None:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                shape.text_frame.clear()

    def _render_title_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Untitled"), top=0.9, size=28)
        self._add_textbox(slide, 0.95, 1.9, 10.5, 0.8, slide_data.get("objective", ""), 14, italic=True)
        for index, bullet in enumerate(slide_data.get("bullets", [])[:3]):
            self._add_textbox(slide, 1.0, 2.7 + index * 0.45, 9.0, 0.35, bullet, 14)

    def _render_agenda_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Agenda"), top=0.6, size=24)
        for index, bullet in enumerate(slide_data.get("bullets", [])[:3]):
            self._add_card(slide, [0.9, 4.45, 8.0][index], 1.7, 2.8, 2.3, f"0{index + 1}", bullet)

    def _render_exec_summary_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Executive Summary"), top=0.65, size=24)
        for index, bullet in enumerate(slide_data.get("bullets", [])[:3]):
            self._add_summary_panel(slide, [0.9, 4.5, 8.1][index], 1.6, 2.9, 3.5, f"Takeaway {index + 1}", bullet)

    def _render_framework_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Framework"), top=0.65, size=24)
        self._add_text_block(slide, 0.85, 1.65, 5.7, 4.8, slide_data.get("bullets", []), 16)
        self._add_callout(slide, 7.1, 1.8, 4.7, 3.6, slide_data.get("visualSuggestion", ""))

    def _render_funnel_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Patient Funnel"), top=0.65, size=24)
        bullets = slide_data.get("bullets", [])[:5]
        for index, bullet in enumerate(bullets):
            shape = slide.shapes.add_shape(
                self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                self.Inches(0.9 + index * 2.35),
                self.Inches(2.0 + index * 0.08),
                self.Inches(2.1),
                self.Inches(3.0 - index * 0.22 if len(bullets) > 3 else 2.5),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.RGBColor(230 - index * 10, 239 - index * 6, 252 - index * 6)
            shape.line.color.rgb = self.RGBColor(30, 64, 175)
            self._set_shape_text(shape, bullet, 13)
        self._add_textbox(slide, 0.95, 5.7, 10.6, 0.45, slide_data.get("visualSuggestion", ""), 12, color=(75, 85, 99), italic=True)

    def _render_insight_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Insight"), top=0.65, size=24)
        self._add_text_block(slide, 0.85, 1.65, 5.3, 4.8, slide_data.get("bullets", []), 16)
        self._add_visual_placeholder(slide, 6.7, 1.65, 5.2, 4.2, slide_data.get("visualSuggestion", ""))

    def _render_recommendation_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Recommendations"), top=0.65, size=24)
        for index, bullet in enumerate(slide_data.get("bullets", [])[:3]):
            self._add_summary_panel(slide, 0.95, 1.7 + index * 1.4, 10.9, 1.0, ["Priority 1", "Priority 2", "Priority 3"][index], bullet)

    def _render_appendix_slide(self, slide, slide_data: dict) -> None:
        self._set_title(slide, slide_data.get("title", "Appendix"), top=0.65, size=24)
        self._add_text_block(slide, 0.95, 1.7, 10.8, 4.9, slide_data.get("bullets", []), 15)

    def _set_title(self, slide, text: str, top: float, size: int) -> None:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type == self.PP_PLACEHOLDER.TITLE:
                    shape.text = text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = self.Pt(size)
                        paragraph.font.bold = True
                    return
            except Exception:
                continue
        self._add_textbox(slide, 0.85, top, 10.5, 0.7, text, size, bold=True)

    def _add_text_block(self, slide, left: float, top: float, width: float, height: float, bullets: list[str], size: int) -> None:
        box = slide.shapes.add_textbox(self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = self.Pt(size)
            paragraph.font.color.rgb = self.RGBColor(31, 41, 55)
            paragraph.space_after = self.Pt(10)

    def _add_visual_placeholder(self, slide, left: float, top: float, width: float, height: float, hint: str) -> None:
        shape = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.RGBColor(255, 255, 255)
        shape.line.color.rgb = self.RGBColor(191, 219, 254)
        default_hint = "Insert the specific chart, network, ranked table, or patient-flow visual that proves the slide headline."
        self._set_shape_text(shape, f"Required visual for this slide\n\n{hint or default_hint}", 14)

    def _add_callout(self, slide, left: float, top: float, width: float, height: float, text: str) -> None:
        shape = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.RGBColor(239, 246, 255)
        shape.line.color.rgb = self.RGBColor(30, 64, 175)
        self._set_shape_text(shape, text or "Add implication, definition, or supporting visual note here.", 14)

    def _add_card(self, slide, left: float, top: float, width: float, height: float, label: str, body: str) -> None:
        shape = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.RGBColor(250, 250, 252)
        shape.line.color.rgb = self.RGBColor(203, 213, 225)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = self.MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = self.Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = self.RGBColor(30, 64, 175)
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = self.Pt(14)
        p2.font.color.rgb = self.RGBColor(31, 41, 55)

    def _add_summary_panel(self, slide, left: float, top: float, width: float, height: float, heading: str, body: str) -> None:
        shape = slide.shapes.add_shape(self.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.RGBColor(255, 255, 255)
        shape.line.color.rgb = self.RGBColor(203, 213, 225)
        tf = shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = heading
        p1.font.size = self.Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = self.RGBColor(30, 64, 175)
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = self.Pt(14)
        p2.font.color.rgb = self.RGBColor(31, 41, 55)

    def _add_textbox(self, slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False, italic: bool = False, color=(31, 41, 55)) -> None:
        box = slide.shapes.add_textbox(self.Inches(left), self.Inches(top), self.Inches(width), self.Inches(height))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = self.Pt(size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.color.rgb = self.RGBColor(*color)

    def _set_shape_text(self, shape, text: str, size: int) -> None:
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = self.Pt(size)
        paragraph.font.color.rgb = self.RGBColor(31, 41, 55)
        paragraph.alignment = self.PP_ALIGN.LEFT

    def _add_footer(self, slide, slide_data: dict) -> None:
        footer = slide.shapes.add_textbox(self.Inches(0.7), self.Inches(6.75), self.Inches(12.0), self.Inches(0.3))
        paragraph = footer.text_frame.paragraphs[0]
        paragraph.text = f"Claim sources: {', '.join(slide_data.get('claimSources', [])) or 'None'} | Style sources: {', '.join(slide_data.get('styleSources', [])) or 'None'}"
        paragraph.font.size = self.Pt(9)
        paragraph.font.color.rgb = self.RGBColor(107, 114, 128)

    def _upsert_footer(self, slide, slide_data: dict) -> None:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and self._shape_text(shape).strip().startswith("Claim sources:"):
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.text = f"Claim sources: {', '.join(slide_data.get('claimSources', [])) or 'None'} | Style sources: {', '.join(slide_data.get('styleSources', [])) or 'None'}"
                paragraph.font.size = self.Pt(9)
                paragraph.font.color.rgb = self.RGBColor(107, 114, 128)
                return
        self._add_footer(slide, slide_data)

    def _attach_notes(self, slide, slide_data: dict) -> None:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        lines = [
            f"Objective: {slide_data.get('objective', '')}",
            f"Visual suggestion: {slide_data.get('visualSuggestion', '')}",
        ]
        if slide_data.get("claimSources"):
            lines.append("Claim sources: " + ", ".join(slide_data.get("claimSources", [])))
        if slide_data.get("styleSources"):
            lines.append("Style sources: " + ", ".join(slide_data.get("styleSources", [])))
        lines.extend(slide_data.get("notes", []))
        notes.text = "\n".join(lines)
