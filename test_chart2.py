from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])

chart_data = CategoryChartData()
chart_data.categories = ['A']
chart_data.add_series('S', (1,))

for shape in slide.placeholders:
    if shape.placeholder_format.type == 18:
        try:
            shape.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
            print("insert_chart worked on PICTURE")
        except Exception as e:
            print("Failed:", e)

prs.save("test_chart.pptx")
